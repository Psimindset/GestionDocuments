import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from PIL import Image as PILImage

OUTPUT_DIR = 'outputs'

def _assurer_dossier():
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)


# Conversion Word
def convertir_word_pdf(chemin_docx: str) -> str:

    _assurer_dossier()

    # Nom du fichier de sortie
    nom_base   = os.path.splitext(os.path.basename(chemin_docx))[0]
    chemin_pdf = os.path.join(OUTPUT_DIR, f'{nom_base}.pdf')

    # Lire le contenu Word
    doc        = Document(chemin_docx)
    styles     = getSampleStyleSheet()
    contenu    = []

    for para in doc.paragraphs:
        if para.text.strip():

            if 'Heading' in para.style.name:
                contenu.append(Paragraph(para.text, styles['Heading1']))
            else:
                contenu.append(Paragraph(para.text, styles['Normal']))
            contenu.append(Spacer(1, 8))

    # Générer le PDF
    pdf = SimpleDocTemplate(chemin_pdf, pagesize=letter)
    pdf.build(contenu)

    print(f'Word → PDF : {chemin_pdf}')
    return chemin_pdf


#  Conversion PowerPoint pptx a PDF
def convertir_pptx_pdf(chemin_pptx: str) -> str:
    _assurer_dossier()

    nom_base   = os.path.splitext(os.path.basename(chemin_pptx))[0]
    chemin_pdf = os.path.join(OUTPUT_DIR, f'{nom_base}.pdf')

    # Lire le contenu PowerPoint
    prs     = Presentation(chemin_pptx)
    styles  = getSampleStyleSheet()
    contenu = []

    for i, slide in enumerate(prs.slides, start=1):
        # Titre de diapositive
        contenu.append(Paragraph(f'── Diapositive {i} ──', styles['Heading2']))
        contenu.append(Spacer(1, 6))

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texte = para.text.strip()
                    if texte:
                        contenu.append(Paragraph(texte, styles['Normal']))
                        contenu.append(Spacer(1, 4))

        contenu.append(Spacer(1, 16))

    # Générer le PDF
    pdf = SimpleDocTemplate(chemin_pdf, pagesize=letter)
    pdf.build(contenu)

    print(f'PPTX → PDF : {chemin_pdf}')
    return chemin_pdf


# Test local
if __name__ == '__main__':
    choix = input('Convertir (1) Word  (2) PowerPoint : ')
    if choix == '1':
        chemin = input('Chemin du fichier .docx : ')
        print(convertir_word_pdf(chemin))
    elif choix == '2':
        chemin = input('Chemin du fichier .pptx : ')
        print(convertir_pptx_pdf(chemin))